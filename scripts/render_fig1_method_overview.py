#!/usr/bin/env python3
# 渲染 Fig 1 Method Overview(可编辑 PPTX)— forest_fire_rgbt
# 路线:python-pptx 手绘矢量 method-figure(非 AI 位图 / 非 matplotlib),每元素独立 shape 可编辑。
# 权威蓝图:04_figures/fig1_method_overview_blueprint.md(3-panel / palette / 标注 / 红线 逐条照做)。
# τ 真值:tau_overlap=0.7 / tau_dual=0.05(01_protocol/formal_protocol_card_20260612.md:125),不编别的数。
# 红线:
#   - Panel B GateA = 决策规则(菱形 + box 操作 + ⊕ add-only),严禁画 neural network / feature concat / attention。
#   - Panel C = schematic(分布区 + 漂移箭头),禁画带数字的真数据曲线 / bar。
#   - 检测框是 schematic 小矩形;真火灾帧用灰块占位标 [RGB-T frame],不嵌真帧。
# 输出:figure_outputs/fig1.pptx(再用 soffice 转 pdf / pdftoppm 转 png 目视)。

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "figure_outputs")
os.makedirs(OUT, exist_ok=True)

# ---------------- palette ----------------
BLUE   = RGBColor(0x4C, 0x72, 0xB0)   # RGB / 视觉模态
RED    = RGBColor(0xC4, 0x4E, 0x52)   # thermal
GREEN  = RGBColor(0x55, 0xA8, 0x68)   # admitted dual
GRAY   = RGBColor(0x8C, 0x8C, 0x8C)   # rejected
INK    = RGBColor(0x33, 0x33, 0x33)   # 主文字
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
# pastel 填充
BLUE_BG  = RGBColor(0xDB, 0xE3, 0xF2)
RED_BG   = RGBColor(0xF3, 0xDD, 0xDE)
GREEN_BG = RGBColor(0xDD, 0xEC, 0xE2)
GRAY_BG  = RGBColor(0xE9, 0xE9, 0xE9)
PANEL_BG = RGBColor(0xF6, 0xF7, 0xF9)   # panel 浅灰底
DIAMOND_BG = RGBColor(0xF5, 0xEF, 0xDC) # 判据米黄
SAFE_BG  = RGBColor(0xDD, 0xEC, 0xE2)   # safe 区(绿)
FRAME_GRAY = RGBColor(0xCF, 0xCF, 0xCF) # 占位帧灰块
BROWN = RGBColor(0x9C, 0x6B, 0x30)

FONT = "Arial"

prs = Presentation()
prs.slide_width  = Inches(12.0)
prs.slide_height = Inches(4.8)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
shapes = slide.shapes


# ---------------- helpers ----------------
def _set_font(run, size, color=INK, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.name = FONT
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.language_id = MSO_LANGUAGE_ID.ENGLISH_US


def textbox(x, y, w, h, lines, size=8, color=INK, bold=False, italic=False,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    tb = shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(1); tf.margin_right = Pt(1)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        txt, kw = (ln if isinstance(ln, tuple) else (ln, {}))
        r = p.add_run(); r.text = txt
        _set_font(r, kw.get("size", size), kw.get("color", color),
                  kw.get("bold", bold), kw.get("italic", italic))
    return tb


def _fill_line(shape, fill, line_color, line_w=1.0, dash=None, no_fill=False):
    if no_fill:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
        if dash:
            ln = shape.line._get_or_add_ln()
            d = ln.find(qn("a:prstDash"))
            if d is None:
                d = ln.makeelement(qn("a:prstDash"), {}); ln.append(d)
            d.set("val", dash)
    shape.shadow.inherit = False


def box(x, y, w, h, fill=WHITE, line=INK, line_w=1.0, dash=None, no_fill=False,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill_line(sp, fill, line, line_w, dash, no_fill)
    return sp


def _put_text(sp, text, size, tcolor, bold, italic=False):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(text, str):
        text = [text]
    for i, ln in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        _set_font(r, size, tcolor, bold, italic)


def labeled_box(x, y, w, h, text, fill=WHITE, line=INK, line_w=1.0, dash=None,
                size=8, tcolor=INK, bold=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                italic=False):
    sp = box(x, y, w, h, fill, line, line_w, dash, shape_type=shape_type)
    _put_text(sp, text, size, tcolor, bold, italic)
    return sp


def diamond(x, y, w, h, text, fill=DIAMOND_BG, line=INK, line_w=1.25, size=7.5, bold=True):
    sp = shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill_line(sp, fill, line, line_w)
    _put_text(sp, text, size, INK, bold)
    return sp


def connector(x1, y1, x2, y2, color=INK, w=1.25, dash=None, arrow=True, begin_arrow=False):
    cn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {}); d.set("val", dash); ln.append(d)
    if arrow:
        tail = ln.makeelement(qn("a:tailEnd"), {})
        tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")
        ln.append(tail)
    if begin_arrow:
        head = ln.makeelement(qn("a:headEnd"), {})
        head.set("type", "triangle"); head.set("w", "med"); head.set("len", "med")
        ln.append(head)
    return cn


def detbox(x, y, w, h, color, dash=None, lw=1.5):
    """schematic 检测框:无填充小矩形,彩色边。"""
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill_line(sp, None, color, lw, dash, no_fill=True)
    return sp


def frame_tile(x, y, w, h, caption, tint=FRAME_GRAY):
    """灰块占位帧 [RGB-T frame],不嵌真帧(schematic 红线)。"""
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill_line(sp, tint, GRAY, 0.75)
    if caption:
        _put_text(sp, caption, 6.5, RGBColor(0x66, 0x66, 0x66), False)
    return sp


def dot(cx, cy, r, fill, line):
    sp = shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r),
                          Inches(2 * r), Inches(2 * r))
    _fill_line(sp, fill, line, 0.75)
    return sp


def panel(x, y, w, h, title, accent):
    bg = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(x), Inches(y), Inches(w), Inches(h))
    _fill_line(bg, PANEL_BG, RGBColor(0xD0, 0xD4, 0xDA), 1.0)
    textbox(x + 0.08, y + 0.04, w - 0.16, 0.30, title, size=10, color=accent,
            bold=True, align=PP_ALIGN.LEFT)
    return bg


# ===================================================================
# 3 panel 横向,B 最大(核心)。预算(12in 宽):
#   A: 0.15 .. 2.70  (w=2.55)
#   B: 2.85 .. 8.80  (w=5.95) 核心最宽
#   C: 8.95 .. 11.90 (w=2.95)
# ===================================================================
PY = 0.62
PH = 4.05
AX, AW = 0.15, 2.55
BX, BW = 2.85, 5.95
CX, CW = 8.95, 2.95

# 大标题
textbox(0.15, 0.06, 11.7, 0.42,
        [("GateA: prediction-level add-only safe admission under LOCO cross-dataset shift",
          {"size": 11.5, "bold": True, "color": INK})],
        align=PP_ALIGN.LEFT)

# =================== Panel A ===================
panel(AX, PY, AW, PH, "(a) LOCO setting", BLUE)

src_names = ["FireMan", "RGBT-3M", "JAG"]
tile_w, tile_h = 0.66, 0.52
tile_y = PY + 0.46
sx0 = AX + 0.14
gap = 0.135
textbox(AX + 0.08, tile_y - 0.205, AW - 0.16, 0.18,
        "source datasets (train)", size=7.5, color=BLUE, bold=True, align=PP_ALIGN.LEFT)
for i, nm in enumerate(src_names):
    tx = sx0 + i * (tile_w + gap)
    frame_tile(tx, tile_y, tile_w, tile_h, "[RGB-T\nframe]")
    textbox(tx - 0.02, tile_y + tile_h + 0.005, tile_w + 0.04, 0.18, nm,
            size=7, color=INK, bold=True)

arr_x = AX + AW / 2.0
model_y = tile_y + tile_h + 0.42
labeled_box(arr_x - 0.52, model_y, 1.04, 0.38, "train", fill=BLUE_BG, line=BLUE,
            line_w=1.25, size=8.5, tcolor=BLUE, bold=True)
for i in range(3):
    tx = sx0 + i * (tile_w + gap) + tile_w / 2.0
    connector(tx, tile_y + tile_h + 0.215, arr_x, model_y - 0.01, color=BLUE, w=1.0)

mb_y = model_y + 0.54
labeled_box(arr_x - 0.60, mb_y, 1.20, 0.42, ["detector", "(RGB / RGB-T)"],
            fill=WHITE, line=INK, line_w=1.25, size=8, bold=True)
connector(arr_x, model_y + 0.38, arr_x, mb_y - 0.01, color=INK, w=1.25)

ho_y = mb_y + 0.62
ho_x = AX + 0.18
frame_tile(ho_x, ho_y, 0.86, 0.56, "[held-out\nRGB-T\nframe]", tint=RGBColor(0xDA, 0xD2, 0xC6))
textbox(ho_x - 0.04, ho_y + 0.56 + 0.005, 0.94, 0.18, "held-out (test)",
        size=7, color=INK, bold=True)
connector(arr_x - 0.42, mb_y + 0.42, ho_x + 0.44, ho_y - 0.02,
          color=RED, w=1.5, dash="dash")
textbox(ho_x + 0.96, ho_y - 0.04, 1.45, 0.66,
        [("distribution shift", {"size": 7, "bold": True, "color": RED}),
         ("Leave-one-dataset-out (LOCO) cross-dataset shift", {"size": 6.8, "color": INK})],
        align=PP_ALIGN.LEFT)

# =================== Panel B (核心) ===================
panel(BX, PY, BW, PH, "(b) GateA add-only admission rule  (prediction-level RULE, NOT a network)", RED)

# --- 输入帧 + 两路 detector ---
in_x = BX + 0.16
in_y = PY + 0.92
frame_tile(in_x, in_y, 0.70, 0.66, "[input\nRGB-T\nframe]")
textbox(in_x - 0.04, in_y + 0.66 + 0.01, 0.78, 0.15, "same frame", size=6.8, color=INK)

prgb_x = BX + 1.08
prgb_y = PY + 0.50
labeled_box(prgb_x, prgb_y, 1.12, 0.44, ["P_rgb", "RGB-only detector"],
            fill=BLUE_BG, line=BLUE, line_w=1.25, size=7.6, tcolor=BLUE, bold=True)
det_y1 = prgb_y + 0.50
detbox(prgb_x + 0.16, det_y1, 0.38, 0.24, BLUE, lw=1.5)
detbox(prgb_x + 0.62, det_y1 + 0.04, 0.32, 0.18, BLUE, lw=1.5)
textbox(prgb_x - 0.05, det_y1 + 0.28, 1.22, 0.15, "blue = RGB boxes (safe default)",
        size=6.5, color=BLUE)

pdual_y = prgb_y + 1.08
labeled_box(prgb_x, pdual_y, 1.12, 0.44, ["P_dual", "RGB-T detector"],
            fill=RED_BG, line=RED, line_w=1.25, size=7.6, tcolor=RED, bold=True)
det_y2 = pdual_y + 0.50
detbox(prgb_x + 0.16, det_y2, 0.38, 0.24, RED, lw=1.5)
detbox(prgb_x + 0.62, det_y2 + 0.04, 0.32, 0.18, RED, lw=1.5)
textbox(prgb_x - 0.05, det_y2 + 0.28, 1.22, 0.15, "red = thermal candidate boxes",
        size=6.5, color=RED)

connector(in_x + 0.70, in_y + 0.16, prgb_x - 0.01, prgb_y + 0.22, color=BLUE, w=1.0)
connector(in_x + 0.70, in_y + 0.50, prgb_x - 0.01, pdual_y + 0.22, color=RED, w=1.0)

# --- 决策菱形 1:IoU >= 0.7 ? ---
d1x = prgb_x + 1.30
d1y = PY + 0.92
dW, dH = 1.02, 0.90
textbox(d1x - 0.10, d1y - 0.215, dW + 0.20, 0.18,
        [("τ_overlap = 0.7  (fixed)", {"size": 6.8, "bold": True, "color": INK})])
diamond(d1x, d1y, dW, dH, ["max IoU vs", "same-class", "RGB box ≥ 0.7 ?"], size=6.6)
connector(prgb_x + 1.12, pdual_y + 0.22, d1x - 0.01, d1y + dH / 2.0, color=RED, w=1.25)

# reject 分支(灰虚线框)向下
rej_x = d1x + 0.02
rej_y = d1y + dH + 0.16
labeled_box(rej_x, rej_y, 1.06, 0.38, ["reject", "(duplicate of RGB)"],
            fill=GRAY_BG, line=GRAY, line_w=1.0, dash="dash", size=6.8, tcolor=GRAY, bold=True)
connector(d1x + dW / 2.0, d1y + dH, rej_x + 0.53, rej_y - 0.01, color=GRAY, w=1.0, dash="dash")
textbox(d1x + dW / 2.0 + 0.03, d1y + dH + 0.005, 0.42, 0.14, "yes",
        size=6.8, color=GRAY, bold=True, align=PP_ALIGN.LEFT)
detbox(rej_x + 0.32, rej_y + 0.42, 0.40, 0.18, GRAY, dash="dash", lw=1.25)

# --- 决策菱形 2:conf >= 0.05 ? ---
d2x = d1x + dW + 0.42
d2y = d1y
textbox(d2x - 0.10, d2y - 0.215, dW + 0.20, 0.18,
        [("τ_dual = 0.05  (fixed)", {"size": 6.8, "bold": True, "color": INK})])
diamond(d2x, d2y, dW, dH, ["conf(d)", "≥ 0.05 ?"], size=7.2)
connector(d1x + dW, d1y + dH / 2.0, d2x - 0.01, d2y + dH / 2.0, color=INK, w=1.25)
textbox(d1x + dW + 0.01, d1y + dH / 2.0 - 0.20, 0.46, 0.14, "no",
        size=6.8, color=INK, bold=True, align=PP_ALIGN.CENTER)

# admit 分支(绿实线框,新增)+ ⊕(add 符号)
adm_x = d2x + 0.02
adm_y = d2y + dH + 0.16
# ⊕ 放在菱形与 admit 框之间,"yes" 标签放右侧避免压住符号
oval = shapes.add_shape(MSO_SHAPE.OVAL, Inches(d2x + dW / 2.0 - 0.15),
                        Inches(d2y + dH + 0.00), Inches(0.30), Inches(0.30))
_fill_line(oval, GREEN_BG, GREEN, 1.5)
plus = shapes.add_shape(MSO_SHAPE.MATH_PLUS, Inches(d2x + dW / 2.0 - 0.085),
                        Inches(d2y + dH + 0.065), Inches(0.17), Inches(0.17))
_fill_line(plus, GREEN, GREEN, 0.75)
textbox(d2x + dW / 2.0 + 0.17, d2y + dH + 0.02, 0.40, 0.14, "yes",
        size=6.8, color=GREEN, bold=True, align=PP_ALIGN.LEFT)
labeled_box(adm_x, adm_y + 0.36, 1.06, 0.38, ["admit (ADD)", "new green box"],
            fill=GREEN_BG, line=GREEN, line_w=1.5, size=6.8, tcolor=GREEN, bold=True)
connector(d2x + dW / 2.0, d2y + dH + 0.30, adm_x + 0.53, adm_y + 0.35, color=GREEN, w=1.5)
detbox(adm_x + 0.33, adm_y + 0.80, 0.40, 0.20, GREEN, lw=1.75)

# --- class-aware NMS + P_out(右侧竖向终端,在 panel 内)---
nms_w = 1.06
term_cx = BX + BW - 0.62           # 终端列中心(贴 panel 右内缘,留 0.09 边距)
nms_x = term_cx - nms_w / 2.0
nms_y = PY + 0.58
labeled_box(nms_x, nms_y, nms_w, 0.42, ["class-aware", "NMS"], fill=WHITE, line=INK,
            line_w=1.25, size=7.4, bold=True)
textbox(term_cx - 0.66, nms_y - 0.205, 1.32, 0.16, "RGB kept as-is (default)",
        size=6.5, color=BLUE, align=PP_ALIGN.CENTER)
# RGB(蓝,默认)-> NMS(从 P_rgb 顶部走上方到 NMS 左上)
connector(prgb_x + 1.12, prgb_y + 0.10, nms_x + 0.02, nms_y + 0.12, color=BLUE, w=1.25)
# admit(绿)-> NMS(从 admit 框右侧到 NMS 左下,绕过 P_out)
connector(adm_x + 1.06, adm_y + 0.50, nms_x + 0.02, nms_y + 0.32, color=GREEN, w=1.25)

# P_out 终端 tile(RGB blue + admitted green)
pout_x = term_cx - 0.58
pout_y = nms_y + 0.60
box(pout_x, pout_y, 1.16, 0.94, fill=WHITE, line=INK, line_w=1.25,
    shape_type=MSO_SHAPE.RECTANGLE)
textbox(pout_x, pout_y + 0.02, 1.16, 0.16, "P_out", size=8, color=INK, bold=True)
fr_x, fr_y = pout_x + 0.11, pout_y + 0.20
frame_tile(fr_x, fr_y, 0.94, 0.62, "")
detbox(fr_x + 0.10, fr_y + 0.10, 0.32, 0.22, BLUE, lw=1.5)
detbox(fr_x + 0.50, fr_y + 0.16, 0.28, 0.18, BLUE, lw=1.5)
detbox(fr_x + 0.20, fr_y + 0.34, 0.30, 0.20, GREEN, lw=1.75)
textbox(pout_x - 0.12, pout_y + 0.94 + 0.005, 1.40, 0.15, "RGB + admitted dual",
        size=6.5, color=INK)
connector(term_cx, nms_y + 0.42, term_cx, pout_y - 0.01, color=INK, w=1.25)

# add-only 强调标注(panel B 底部)
textbox(BX + 0.15, PY + PH - 0.48, BW - 0.30, 0.44,
        [("add-only: thermal detections are ADDED (⊕), never replace or suppress RGB",
          {"size": 8, "bold": True, "color": GREEN}),
         ("default output = RGB-only (safe);  τ_overlap = 0.7 and τ_dual = 0.05 are fixed, not learned",
          {"size": 7, "color": INK})],
        align=PP_ALIGN.LEFT)

# =================== Panel C (schematic, 禁真数据) ===================
panel(CX, PY, CW, PH, "(c) why it transfers (C014)", GREEN)

zone_w = CW - 0.36
zx = CX + 0.18
zone_h = 1.16
src_zone_y = PY + 0.52
box(zx, src_zone_y, zone_w, zone_h, fill=RGBColor(0xE7, 0xEC, 0xF4),
    line=RGBColor(0xB7, 0xC4, 0xDC), line_w=1.0)
textbox(zx + 0.06, src_zone_y + 0.02, zone_w - 0.12, 0.16, "source distribution",
        size=7, color=BLUE, bold=True, align=PP_ALIGN.LEFT)
tgt_zone_y = src_zone_y + zone_h + 0.22
box(zx, tgt_zone_y, zone_w, zone_h, fill=RGBColor(0xF1, 0xEA, 0xE2),
    line=RGBColor(0xD9, 0xC7, 0xB4), line_w=1.0)
textbox(zx + 0.06, tgt_zone_y + 0.02, zone_w - 0.12, 0.16, "target distribution (held-out)",
        size=7, color=BROWN, bold=True, align=PP_ALIGN.LEFT)

# safe band(绿虚线)在每个 zone 下半
band_x = zx + 0.12
band_w = zone_w - 0.24
sb_h = 0.42
sb1_y = src_zone_y + 0.66
sb2_y = tgt_zone_y + 0.66
box(band_x, sb1_y, band_w, sb_h, fill=SAFE_BG, line=GREEN, line_w=0.75, dash="dash",
    shape_type=MSO_SHAPE.RECTANGLE)
# 源 zone 的 "safe region" 标签放底部居中,避开顶部 learned-gate / GateA 标签
textbox(band_x + 0.04, sb1_y + sb_h - 0.155, band_w - 0.08, 0.14, "safe region",
        size=6.5, color=GREEN, align=PP_ALIGN.CENTER)
box(band_x, sb2_y, band_w, sb_h, fill=SAFE_BG, line=GREEN, line_w=0.75, dash="dash",
    shape_type=MSO_SHAPE.RECTANGLE)
textbox(band_x + 0.04, sb2_y + 0.01, 0.9, 0.14, "safe region", size=6.5, color=GREEN, align=PP_ALIGN.LEFT)

# learned gate:学在 source(safe 内 左侧),漂向 target(掉出 safe -> 落 target 上半 unsafe)
lg_x = band_x + 0.34
lg_src_y = sb1_y + sb_h / 2.0
dot(lg_x, lg_src_y, 0.065, RED, RED)
# 标签放红点左上,左对齐贴左缘,避开右侧 GateA 标签
textbox(lg_x - 0.30, lg_src_y - 0.235, 1.45, 0.15, "learned gate",
        size=6.2, color=RED, bold=True, align=PP_ALIGN.LEFT)
lg_tgt_y = tgt_zone_y + 0.40   # target safe band 之上 = 掉出 safe(下移避开 zone 标题)
connector(lg_x, lg_src_y + 0.07, lg_x, lg_tgt_y, color=RED, w=1.75, dash="dash")
dot(lg_x, lg_tgt_y, 0.065, RED, RED)
textbox(lg_x + 0.12, lg_tgt_y - 0.075, 1.55, 0.18, "misfit · negative transfer",
        size=6.2, color=RED, bold=True, align=PP_ALIGN.LEFT)

# GateA 固定点:τ=0.7/0.05,在 source/target 两区都落 safe,不动(垂直实线)
ga_x = band_x + band_w - 0.32
ga_src_y = sb1_y + sb_h / 2.0
ga_tgt_y = sb2_y + sb_h / 2.0
dot(ga_x, ga_src_y, 0.075, GREEN, GREEN)
dot(ga_x, ga_tgt_y, 0.075, GREEN, GREEN)
connector(ga_x, ga_src_y + 0.075, ga_x, ga_tgt_y - 0.075, color=GREEN, w=1.75)
# GateA 标签放绿点右上,右对齐贴右缘,避开左侧 learned-gate 标签
textbox(ga_x - 1.35, ga_src_y - 0.235, 1.30, 0.15, "GateA  τ = 0.7 / 0.05",
        size=6.2, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)
textbox(ga_x - 1.55, ga_tgt_y + 0.085, 1.50, 0.15, "distribution-agnostic · stays safe",
        size=6.2, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)

# Panel C 底部解读句
textbox(CX + 0.15, PY + PH - 0.72, CW - 0.30, 0.68,
        [("A learned gate's operating point does not transfer across datasets;",
          {"size": 7, "color": INK}),
         ("GateA's fixed threshold is the distribution-agnostic safety floor (C014).",
          {"size": 7, "bold": True, "color": INK})],
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)

# =================== panel 间流向箭头(A->B->C)===================
connector(AX + AW + 0.01, PY + PH / 2.0, BX - 0.01, PY + PH / 2.0, color=INK, w=2.0)
connector(BX + BW + 0.01, PY + PH / 2.0, CX - 0.01, PY + PH / 2.0, color=INK, w=2.0)

# ---------------- 保存 ----------------
out_pptx = os.path.join(OUT, "fig1.pptx")
prs.save(out_pptx)
print("OK ->", out_pptx)
